"""
Document ingestion and processing.

This module handles ingesting documents (PDF, doc, etc.) into Gianluigi's
knowledge base:
1. Extract text from various formats
2. Generate summary using Claude
3. Create searchable embeddings
4. Store in Supabase

Supported formats:
- PDF (.pdf)
- Word documents (.docx)
- PowerPoint presentations (.pptx)
- Plain text (.txt)
- Markdown (.md)
- Google Docs (exported as plain text via Drive API)

Usage:
    from processors.document_processor import process_document

    result = await process_document(
        content="...",
        title="Competitor Analysis",
        source="drive",
        file_type="pdf"
    )
"""

import io
import json
import logging
from typing import Any

from config.settings import settings
from core.llm import call_llm
from services.supabase_client import supabase_client
from services.embeddings import embedding_service

# Valid document types for classification
DOCUMENT_TYPES = ["strategy", "legal", "technical", "pitch", "client", "other"]

logger = logging.getLogger(__name__)


async def process_document(
    content: str,
    title: str,
    source: str,
    file_type: str | None = None,
    drive_path: str | None = None
) -> dict:
    """
    Process a document for ingestion into the knowledge base.

    Full pipeline:
    1. Generate summary using Claude
    2. Store document record in Supabase
    3. Chunk text and generate embeddings
    4. Store embeddings in pgvector

    Args:
        content: The document text content (already extracted from PDF/docx).
        title: Document title.
        source: Where it came from ('upload', 'email', 'drive').
        file_type: File extension (pdf, docx, txt, md).
        drive_path: Google Drive path if applicable.

    Returns:
        Dict containing:
        - document_id: UUID of the created document
        - title: Document title
        - summary: Generated summary
        - chunk_count: Number of embedded chunks
    """
    if not content or not content.strip():
        logger.warning(f"Empty content for document: {title}")
        return {
            "document_id": None,
            "title": title,
            "summary": "",
            "chunk_count": 0,
            "error": "Empty document content",
        }

    logger.info(f"Processing document: {title} ({len(content)} chars)")

    # Step 1: Generate summary using Claude
    summary = await generate_document_summary(content, title)

    # Step 2: Classify document type
    document_type = await classify_document_type(content, title)

    # Step 3: Store document record in Supabase
    document = supabase_client.create_document(
        title=title,
        source=source,
        file_type=file_type,
        summary=summary,
        drive_path=drive_path,
        document_type=document_type,
    )
    document_id = document["id"]

    # Step 4: Chunk and embed, then store
    chunk_count = await store_document_embeddings(document_id, content)

    logger.info(
        f"Document processed: {title} — "
        f"ID={document_id}, type={document_type}, chunks={chunk_count}"
    )

    return {
        "document_id": document_id,
        "title": title,
        "summary": summary,
        "document_type": document_type,
        "chunk_count": chunk_count,
    }


async def generate_document_summary(
    content: str,
    title: str,
    max_length: int = 500
) -> str:
    """
    Generate a summary of a document using Claude.

    Args:
        content: The full document text.
        title: Document title for context.
        max_length: Target summary length in words.

    Returns:
        Summary text.
    """
    # Truncate very long documents to stay within token limits.
    # ~4 chars per token, keep within a safe window for Claude.
    max_chars = 80_000
    truncated = content[:max_chars]
    if len(content) > max_chars:
        truncated += "\n\n[... document truncated for summarization ...]"

    prompt = f"""Summarize the following document for CropSight's internal knowledge base.

Document title: {title}

Guidelines:
- Write a concise summary (up to {max_length} words)
- Focus on key facts, decisions, data points, and action-relevant information
- If the document mentions people, organizations, dates, or commitments, include them
- Use professional, factual tone — no opinions or editorializing
- If it's a research paper, include the main findings and methodology
- If it's a contract/MOU, include parties, key terms, and dates

Document content:
{truncated}"""

    try:
        summary, _ = call_llm(
            prompt=prompt,
            model=settings.model_simple,
            max_tokens=1024,
            call_site="document_summary",
        )
        summary = summary.strip()
        logger.info(f"Generated summary for '{title}': {len(summary)} chars")
        return summary

    except Exception as e:
        logger.error(f"Error generating document summary: {e}")
        return f"[Summary generation failed: {e}]"


async def classify_document_type(content: str, title: str) -> str:
    """
    Classify a document into a type category using Claude.

    Types: strategy, legal, technical, pitch, client, other.

    Args:
        content: The document text (first portion used for classification).
        title: Document title for context.

    Returns:
        One of the DOCUMENT_TYPES strings.
    """
    # Use just the first 2000 chars + title for classification (cheap)
    snippet = content[:2000]

    prompt = f"""Classify this document into exactly one category.

Title: {title}

Content preview:
{snippet}

Categories:
- strategy: strategy docs, competitive analyses, market research, business plans
- legal: MOUs, NDAs, contracts, terms & conditions, compliance docs
- technical: research papers, technical specs, API docs, architecture docs
- pitch: pitch decks, investor materials, fundraising docs
- client: client-facing proposals, reports, presentations for external partners
- other: anything that doesn't fit the above

Reply with ONLY the category name (one word, lowercase). Nothing else."""

    try:
        response_text, _ = call_llm(
            prompt=prompt,
            model=settings.model_simple,
            max_tokens=16,
            call_site="document_classification",
        )
        doc_type = response_text.strip().lower()
        if doc_type in DOCUMENT_TYPES:
            return doc_type
        logger.warning(f"Unknown document type '{doc_type}', defaulting to 'other'")
        return "other"
    except Exception as e:
        logger.warning(f"Document classification failed: {e}")
        return "other"


async def extract_key_points(content: str) -> list[str]:
    """
    Extract key points from a document.

    Args:
        content: The document text.

    Returns:
        List of key point strings.
    """
    max_chars = 80_000
    truncated = content[:max_chars]

    prompt = f"""Extract the key points from this document as a numbered list.
Each point should be one concise sentence capturing a distinct fact, decision, or insight.
Return only the numbered list, nothing else.

Document:
{truncated}"""

    try:
        text, _ = call_llm(
            prompt=prompt,
            model=settings.model_simple,
            max_tokens=1024,
            call_site="document_key_points",
        )

        # Parse numbered list into individual points
        points = []
        for line in text.strip().split("\n"):
            line = line.strip()
            if line:
                # Strip leading number + period/parenthesis
                import re
                cleaned = re.sub(r'^\d+[\.\)]\s*', '', line)
                if cleaned:
                    points.append(cleaned)

        return points

    except Exception as e:
        logger.error(f"Error extracting key points: {e}")
        return []


def extract_text_from_pdf(pdf_bytes: bytes) -> str:
    """
    Extract text content from a PDF file.

    Args:
        pdf_bytes: Raw PDF file bytes.

    Returns:
        Extracted text content.
    """
    try:
        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(pdf_bytes))
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)

        content = "\n\n".join(pages)
        logger.info(f"Extracted {len(content)} chars from PDF ({len(reader.pages)} pages)")
        return content

    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        return ""


def extract_text_from_docx(docx_bytes: bytes) -> str:
    """
    Extract text content from a Word document.

    Args:
        docx_bytes: Raw .docx file bytes.

    Returns:
        Extracted text content.
    """
    try:
        from docx import Document

        doc = Document(io.BytesIO(docx_bytes))
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        content = "\n\n".join(paragraphs)
        logger.info(f"Extracted {len(content)} chars from DOCX ({len(paragraphs)} paragraphs)")
        return content

    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {e}")
        return ""


def extract_text_from_pptx(pptx_bytes: bytes) -> str:
    """
    Extract text content from a PowerPoint presentation.

    Preserves slide order as chunk order. Speaker notes are included
    as a separate section per slide.

    Args:
        pptx_bytes: Raw .pptx file bytes.

    Returns:
        Extracted text content with slide structure preserved.
    """
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(pptx_bytes))
        slides = []

        for i, slide in enumerate(prs.slides, 1):
            slide_lines = [f"--- Slide {i} ---"]

            # Extract text from all shapes on the slide
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_lines.append(text)

            # Extract speaker notes if present
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_lines.append(f"[Speaker Notes: {notes_text}]")

            slides.append("\n".join(slide_lines))

        content = "\n\n".join(slides)
        logger.info(f"Extracted {len(content)} chars from PPTX ({len(prs.slides)} slides)")
        return content

    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {e}")
        return ""


def extract_text_by_mime_type(file_bytes: bytes, mime_type: str, filename: str) -> str:
    """
    Extract text from file bytes based on MIME type.

    Handles routing to the correct extraction method.

    Args:
        file_bytes: Raw file bytes.
        mime_type: MIME type from Google Drive.
        filename: Original filename (used as fallback for type detection).

    Returns:
        Extracted text content.
    """
    lower_name = filename.lower()

    # PDF
    if mime_type == "application/pdf" or lower_name.endswith(".pdf"):
        return extract_text_from_pdf(file_bytes)

    # Word documents
    if (
        mime_type in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        )
        or lower_name.endswith(".docx")
        or lower_name.endswith(".doc")
    ):
        return extract_text_from_docx(file_bytes)

    # PowerPoint presentations
    if (
        mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or lower_name.endswith(".pptx")
    ):
        return extract_text_from_pptx(file_bytes)

    # Plain text / markdown — decode as UTF-8
    if (
        mime_type in ("text/plain", "text/markdown", "text/csv")
        or lower_name.endswith((".txt", ".md", ".csv"))
    ):
        return file_bytes.decode("utf-8", errors="ignore")

    # Google Docs are handled by Drive API export (returns text directly),
    # so they won't normally reach here. Log a warning for unknown types.
    logger.warning(
        f"Unknown file type for extraction: mime={mime_type}, name={filename}. "
        f"Attempting UTF-8 decode."
    )
    return file_bytes.decode("utf-8", errors="ignore")


DOCUMENT_KEYWORDS = {
    "deck", "pitch", "presentation", "slides", "pptx",
    "doc", "document", "paper", "report", "analysis",
    "proposal", "contract", "mou", "nda", "agreement",
    "brief", "abstract", "memo", "plan", "strategy",
}


def find_related_documents(
    text: str,
    limit: int = 5,
) -> list[dict]:
    """
    Find documents that might be related to text content (e.g., a transcript).

    Searches for document-related keywords in the text, then matches
    against document titles in the database.

    Args:
        text: Text to search for document references (e.g., transcript).
        limit: Maximum number of documents to return.

    Returns:
        List of matching document records.
    """
    # Extract potential document references from the text
    text_lower = text.lower()
    words = set(text_lower.split())

    # Only proceed if the text mentions document-related keywords
    if not words & DOCUMENT_KEYWORDS:
        return []

    # Get all documents and match by keyword overlap with titles
    try:
        all_docs = supabase_client.list_documents(limit=100)
        if not all_docs:
            return []

        matches = []
        for doc in all_docs:
            doc_title = doc.get("title", "").lower()
            # Check if any significant words from the doc title appear in the text
            title_words = set(doc_title.split()) - {"the", "a", "an", "of", "for", "and", "in", "to"}
            if title_words and len(title_words & words) >= 1:
                matches.append(doc)

        return matches[:limit]

    except Exception as e:
        logger.warning(f"Error finding related documents: {e}")
        return []


async def store_document_embeddings(
    document_id: str,
    content: str
) -> int:
    """
    Chunk document and store embeddings in Supabase pgvector.

    Args:
        document_id: UUID of the document.
        content: Document text content.

    Returns:
        Number of chunks stored.
    """
    if not content or not content.strip():
        return 0

    try:
        # Chunk and embed using the embedding service
        embedded_chunks = await embedding_service.chunk_and_embed_document(
            document=content,
            document_id=document_id,
        )

        if not embedded_chunks:
            logger.warning(f"No chunks generated for document {document_id}")
            return 0

        # Prepare records for batch insert
        embedding_records = [
            {
                "source_type": "document",
                "source_id": document_id,
                "chunk_text": chunk["text"],
                "chunk_index": chunk["chunk_index"],
                "embedding": chunk["embedding"],
                "speaker": None,
                "timestamp_range": None,
                "metadata": chunk.get("metadata", {}),
            }
            for chunk in embedded_chunks
        ]

        supabase_client.store_embeddings_batch(embedding_records)

        logger.info(
            f"Stored {len(embedding_records)} embeddings for document {document_id}"
        )
        return len(embedding_records)

    except Exception as e:
        logger.error(f"Error storing document embeddings: {e}")
        return 0
