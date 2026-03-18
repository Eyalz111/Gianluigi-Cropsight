"""
Gantt slide (PPTX) generator.

Generates a table-based PPTX slide showing the operational Gantt chart
status for the weekly review. Professional layout with colored cells.
"""

import io
import logging
from datetime import datetime, timedelta

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

# Status -> color mapping
STATUS_COLORS = {
    "active": RGBColor(0x4C, 0xAF, 0x50),    # green
    "planned": RGBColor(0x21, 0x96, 0xF3),    # blue
    "blocked": RGBColor(0xF4, 0x43, 0x36),    # red
    "completed": RGBColor(0x9E, 0x9E, 0x9E),  # gray
    "delayed": RGBColor(0xFF, 0x98, 0x00),     # orange
    "": RGBColor(0xF5, 0xF5, 0xF5),            # light gray (empty)
}

# Owner abbreviations
OWNER_LEGEND = {
    "Eyal": "[E]",
    "Roye": "[R]",
    "Paolo": "[P]",
    "Yoram": "[Y]",
    "Prof. Yoram Weiss": "[Y]",
}

# Milestone markers
MILESTONE_MARKERS = {"★": "tech", "●": "commercial", "◆": "funding"}

# Gantt sections
GANTT_SECTIONS = [
    "Strategic Milestones",
    "Product & Tech",
    "Sales & BD",
    "Fundraising",
    "Legal & Finance",
]


async def generate_gantt_slide(
    week_number: int,
    year: int,
    week_range: tuple[int, int] | None = None,
) -> bytes:
    """
    Generate a PPTX slide with Gantt chart table.

    Args:
        week_number: Current ISO week number.
        year: Year.
        week_range: Optional (start_week, end_week) for columns.
                    Defaults to current quarter + 4 weeks.

    Returns:
        PPTX file as bytes.
    """
    # Get Gantt data
    gantt_data = {}
    try:
        from services.gantt_manager import gantt_manager
        gantt_data = await gantt_manager.get_gantt_status()
    except Exception as e:
        logger.warning(f"Could not read Gantt data: {e}")

    # Calculate week range
    if week_range is None:
        start_week = max(1, week_number - 2)
        end_week = min(53, week_number + 10)
        week_range = (start_week, end_week)

    pptx_bytes = _build_slide(gantt_data, week_number, year, week_range)

    logger.info(
        f"PPTX slide generated for W{week_number}/{year}, "
        f"size={len(pptx_bytes)} bytes"
    )
    return pptx_bytes


def _build_slide(
    gantt_data: dict,
    week_number: int,
    year: int,
    week_range: tuple[int, int],
) -> bytes:
    """Build the PPTX presentation with table-based Gantt."""
    prs = Presentation()

    # Widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    _add_title(slide, week_number, year)

    # Build section tables
    sections = gantt_data.get("sections", {})
    y_offset = Inches(1.2)

    for section_name in GANTT_SECTIONS:
        section_data = sections.get(section_name, {})
        items = section_data.get("items", [])

        if not items:
            # Still show section header even if empty
            items = [{"name": "(no items)", "weeks": {}}]

        y_offset = _add_section_table(
            slide, section_name, items,
            week_range, week_number, y_offset,
        )
        y_offset += Inches(0.15)

    # Owner legend
    _add_legend(slide, y_offset)

    # Footer
    _add_footer(slide)

    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.read()


def _add_title(slide, week_number: int, year: int) -> None:
    """Add title text box to slide."""
    from pptx.util import Inches, Pt

    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(10), Inches(0.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = f"CropSight Operational Gantt — Week {week_number}, {year}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2D, 0x50, 0x16)

    # Badge
    p2 = tf.add_paragraph()
    p2.text = "Gianluigi-generated"
    p2.font.size = Pt(10)
    p2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p2.font.italic = True


def _add_section_table(
    slide,
    section_name: str,
    items: list[dict],
    week_range: tuple[int, int],
    current_week: int,
    y_offset,
) -> float:
    """Add a section table to the slide. Returns new y_offset."""
    from pptx.util import Inches, Pt

    start_week, end_week = week_range
    num_weeks = end_week - start_week + 1
    num_cols = 2 + num_weeks  # Item name + Owner + week columns
    num_rows = 1 + len(items)  # Header + items

    # Table dimensions
    table_width = Inches(12.3)
    table_height = Inches(0.25 * num_rows)

    # Cap table height
    max_height = Inches(1.5)
    if table_height > max_height:
        table_height = max_height

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.5), y_offset,
        table_width, table_height,
    )
    table = table_shape.table

    # Header row
    _set_cell(table.cell(0, 0), section_name, bold=True, size=9,
              bg=RGBColor(0x2D, 0x50, 0x16), fg=RGBColor(0xFF, 0xFF, 0xFF))
    _set_cell(table.cell(0, 1), "Owner", bold=True, size=8,
              bg=RGBColor(0x2D, 0x50, 0x16), fg=RGBColor(0xFF, 0xFF, 0xFF))

    for w_idx in range(num_weeks):
        w = start_week + w_idx
        is_current = (w == current_week)
        label = f"W{w}"
        bg = RGBColor(0x4A, 0x7C, 0x23) if is_current else RGBColor(0x2D, 0x50, 0x16)
        _set_cell(table.cell(0, 2 + w_idx), label, bold=is_current, size=7,
                  bg=bg, fg=RGBColor(0xFF, 0xFF, 0xFF))

    # Data rows
    for row_idx, item in enumerate(items, 1):
        name = item.get("name", "")
        owner = item.get("owner", "")
        owner_abbr = OWNER_LEGEND.get(owner, owner[:3] if owner else "")

        # Check for milestone markers
        marker_info = _parse_markers(name)

        _set_cell(table.cell(row_idx, 0), name, size=8)
        _set_cell(table.cell(row_idx, 1), owner_abbr, size=8,
                  align=PP_ALIGN.CENTER)

        weeks_data = item.get("weeks", {})
        for w_idx in range(num_weeks):
            w = start_week + w_idx
            status = weeks_data.get(str(w), weeks_data.get(w, ""))
            color = _status_to_color(status)

            cell_text = ""
            if marker_info["is_milestone"]:
                cell_text = marker_info["marker"] if status else ""

            cell = table.cell(row_idx, 2 + w_idx)
            _set_cell(cell, cell_text, size=7, bg=color if status else None,
                      align=PP_ALIGN.CENTER)

            # Highlight current week column
            if w == current_week and not status:
                _set_cell_bg(cell, RGBColor(0xF0, 0xF7, 0xEB))

    # Set column widths
    col_widths = table.columns
    col_widths[0].width = Inches(3.0)
    col_widths[1].width = Inches(0.5)
    week_col_width = Inches((12.3 - 3.5) / num_weeks)
    for i in range(2, num_cols):
        col_widths[i].width = int(week_col_width)

    return y_offset + table_height


def _add_legend(slide, y_offset) -> None:
    """Add owner legend."""
    from pptx.util import Inches, Pt

    txBox = slide.shapes.add_textbox(
        Inches(0.5), y_offset, Inches(10), Inches(0.3)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = (
        "Owners: [E] Eyal  [R] Roye  [P] Paolo  [Y] Prof. Yoram    "
        "Milestones: ★ Tech  ● Commercial  ◆ Funding    "
        "Colors: Green=Active  Blue=Planned  Red=Blocked  Gray=Completed"
    )
    p.font.size = Pt(8)
    p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)


def _add_footer(slide) -> None:
    """Add confidentiality footer."""
    from pptx.util import Inches, Pt

    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.0), Inches(12), Inches(0.3)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = (
        f"CropSight — Confidential | "
        f"Generated by Gianluigi on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
    )
    p.font.size = Pt(7)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER


# =========================================================================
# Helpers
# =========================================================================

def _status_to_color(status: str) -> RGBColor:
    """Map status string to RGBColor."""
    status_lower = status.lower().strip() if status else ""
    return STATUS_COLORS.get(status_lower, STATUS_COLORS.get("", RGBColor(0xF5, 0xF5, 0xF5)))


def _parse_markers(text: str) -> dict:
    """Parse milestone markers from text."""
    for marker, mtype in MILESTONE_MARKERS.items():
        if marker in text:
            return {"is_milestone": True, "marker": marker, "type": mtype}
    return {"is_milestone": False, "marker": None, "type": None}


def _set_cell(
    cell,
    text: str,
    bold: bool = False,
    size: int = 9,
    bg: RGBColor | None = None,
    fg: RGBColor | None = None,
    align=None,
) -> None:
    """Set cell text and formatting."""
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        if fg:
            paragraph.font.color.rgb = fg
        if align:
            paragraph.alignment = align

    if bg:
        _set_cell_bg(cell, bg)


def _set_cell_bg(cell, color: RGBColor) -> None:
    """Set cell background color."""
    from pptx.oxml.ns import qn
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(
        qn("a:srgbClr"),
        {"val": f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"},
    )
    solidFill.append(srgbClr)
    tcPr.append(solidFill)
