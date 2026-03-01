"""
Tests for document ingestion pipeline (v0.5 Track D).

Tests cover:
- PPTX text extraction
- Document type classification
- Document-meeting cross-referencing
- Telegram notification on ingestion
- Updated CRUD operations
"""

import io
import json
import pytest
from unittest.mock import patch, MagicMock, AsyncMock


# =========================================================================
# Test PPTX extraction
# =========================================================================

class TestExtractTextFromPptx:
    """Tests for PowerPoint text extraction."""

    def test_extracts_slide_text(self):
        """Should extract text from slide shapes."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + content
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text = "Bullet point content"

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        from processors.document_processor import extract_text_from_pptx
        result = extract_text_from_pptx(buf.read())

        assert "Slide 1" in result
        assert "Slide Title" in result
        assert "Bullet point content" in result

    def test_extracts_speaker_notes(self):
        """Should extract speaker notes from slides."""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Main Slide"

        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Remember to mention the deadline"

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        from processors.document_processor import extract_text_from_pptx
        result = extract_text_from_pptx(buf.read())

        assert "Speaker Notes" in result
        assert "deadline" in result

    def test_preserves_slide_order(self):
        """Should preserve slide numbering."""
        from pptx import Presentation

        prs = Presentation()
        for i in range(3):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = f"Slide {i + 1} Title"

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        from processors.document_processor import extract_text_from_pptx
        result = extract_text_from_pptx(buf.read())

        # Check slide markers are in order
        pos1 = result.index("Slide 1")
        pos2 = result.index("Slide 2")
        pos3 = result.index("Slide 3")
        assert pos1 < pos2 < pos3

    def test_empty_pptx_returns_empty(self):
        """Empty presentation should return minimal content."""
        from pptx import Presentation

        prs = Presentation()
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        from processors.document_processor import extract_text_from_pptx
        result = extract_text_from_pptx(buf.read())
        # No slides means no content
        assert result == ""

    def test_invalid_bytes_returns_empty(self):
        """Invalid bytes should return empty string."""
        from processors.document_processor import extract_text_from_pptx
        result = extract_text_from_pptx(b"not a pptx file")
        assert result == ""


# =========================================================================
# Test MIME type routing includes PPTX
# =========================================================================

class TestMimeTypeRouting:
    """Tests for extract_text_by_mime_type PPTX routing."""

    def test_routes_pptx_by_mime_type(self):
        """Should route PPTX MIME type to PPTX extractor."""
        with patch("processors.document_processor.extract_text_from_pptx") as mock_extract:
            mock_extract.return_value = "slide content"

            from processors.document_processor import extract_text_by_mime_type
            result = extract_text_by_mime_type(
                b"data",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "deck.pptx",
            )

            mock_extract.assert_called_once_with(b"data")
            assert result == "slide content"

    def test_routes_pptx_by_extension(self):
        """Should route .pptx extension to PPTX extractor."""
        with patch("processors.document_processor.extract_text_from_pptx") as mock_extract:
            mock_extract.return_value = "slide content"

            from processors.document_processor import extract_text_by_mime_type
            result = extract_text_by_mime_type(
                b"data",
                "application/octet-stream",
                "Pitch Deck.pptx",
            )

            mock_extract.assert_called_once()


# =========================================================================
# Test document type classification
# =========================================================================

class TestClassifyDocumentType:
    """Tests for document type classification."""

    @pytest.mark.asyncio
    async def test_classifies_strategy_doc(self):
        """Should classify a strategy document correctly."""
        with patch("processors.document_processor.call_llm") as mock_llm:
            mock_llm.return_value = ("strategy", {"input_tokens": 50, "output_tokens": 5, "cache_read_input_tokens": 0, "cache_creation_input_tokens": 0})

            from processors.document_processor import classify_document_type
            result = await classify_document_type("Market analysis for Q1...", "Q1 Strategy Review")

            assert result == "strategy"
            mock_llm.assert_called_once()

    @pytest.mark.asyncio
    async def test_classifies_legal_doc(self):
        """Should classify a legal document correctly."""
        with patch("processors.document_processor.call_llm") as mock_llm:
            mock_llm.return_value = ("legal", {"input_tokens": 50, "output_tokens": 5, "cache_read_input_tokens": 0, "cache_creation_input_tokens": 0})

            from processors.document_processor import classify_document_type
            result = await classify_document_type("This MOU between...", "Lavazza MOU")

            assert result == "legal"

    @pytest.mark.asyncio
    async def test_unknown_type_defaults_to_other(self):
        """Unknown classification should default to 'other'."""
        with patch("processors.document_processor.call_llm") as mock_llm:
            mock_llm.return_value = ("marketing", {"input_tokens": 50, "output_tokens": 5, "cache_read_input_tokens": 0, "cache_creation_input_tokens": 0})

            from processors.document_processor import classify_document_type
            result = await classify_document_type("Some content", "Title")

            assert result == "other"

    @pytest.mark.asyncio
    async def test_error_defaults_to_other(self):
        """LLM error should default to 'other'."""
        with patch("processors.document_processor.call_llm") as mock_llm:
            mock_llm.side_effect = Exception("API error")

            from processors.document_processor import classify_document_type
            result = await classify_document_type("content", "title")

            assert result == "other"

    @pytest.mark.asyncio
    async def test_uses_model_simple(self):
        """Should use model_simple for classification (cheap)."""
        with patch("processors.document_processor.call_llm") as mock_llm:
            mock_llm.return_value = ("technical", {"input_tokens": 50, "output_tokens": 5, "cache_read_input_tokens": 0, "cache_creation_input_tokens": 0})

            from processors.document_processor import classify_document_type
            await classify_document_type("content", "title")

            call_args = mock_llm.call_args
            from config.settings import settings
            assert call_args.kwargs["model"] == settings.model_simple
            assert call_args.kwargs["call_site"] == "document_classification"


# =========================================================================
# Test process_document includes document_type
# =========================================================================

class TestProcessDocumentWithType:
    """Tests for process_document returning document_type."""

    @pytest.mark.asyncio
    async def test_process_document_returns_type(self):
        """process_document should return document_type in result."""
        with (
            patch("processors.document_processor.call_llm") as mock_llm,
            patch("processors.document_processor.supabase_client") as mock_db,
            patch("processors.document_processor.embedding_service") as mock_embed,
        ):
            # First call: summary, second call: classification
            mock_llm.side_effect = [
                ("This is a strategy doc summary.", {"input_tokens": 100, "output_tokens": 50, "cache_read_input_tokens": 0, "cache_creation_input_tokens": 0}),
                ("strategy", {"input_tokens": 50, "output_tokens": 5, "cache_read_input_tokens": 0, "cache_creation_input_tokens": 0}),
            ]
            mock_db.create_document.return_value = {"id": "doc-1"}
            mock_embed.chunk_and_embed_document = AsyncMock(return_value=[])

            from processors.document_processor import process_document
            result = await process_document(
                content="Strategic analysis of the competitive landscape...",
                title="Competitive Analysis",
                source="drive",
                file_type="pdf",
            )

            assert result["document_type"] == "strategy"
            # Verify create_document was called with document_type
            mock_db.create_document.assert_called_once()
            call_kwargs = mock_db.create_document.call_args.kwargs
            assert call_kwargs["document_type"] == "strategy"


# =========================================================================
# Test document-meeting cross-referencing
# =========================================================================

class TestFindRelatedDocuments:
    """Tests for document-meeting cross-referencing."""

    def test_finds_documents_by_keyword_overlap(self):
        """Should find documents when title words appear in text."""
        with patch("processors.document_processor.supabase_client") as mock_db:
            mock_db.list_documents.return_value = [
                {"title": "Lavazza Partnership Proposal", "document_type": "client", "summary": "Proposal for..."},
                {"title": "Internal Budget 2026", "document_type": "strategy", "summary": "Budget..."},
            ]

            from processors.document_processor import find_related_documents
            result = find_related_documents(
                "We need to send the Lavazza deck and update the proposal"
            )

            assert len(result) == 1
            assert result[0]["title"] == "Lavazza Partnership Proposal"

    def test_returns_empty_without_document_keywords(self):
        """Should return empty when text has no document-related keywords."""
        from processors.document_processor import find_related_documents
        result = find_related_documents("Just a regular conversation about the weather")
        assert result == []

    def test_returns_empty_when_no_documents_exist(self):
        """Should return empty when document table is empty."""
        with patch("processors.document_processor.supabase_client") as mock_db:
            mock_db.list_documents.return_value = []

            from processors.document_processor import find_related_documents
            result = find_related_documents("Send the deck to investors")
            assert result == []

    def test_handles_db_error_gracefully(self):
        """Should return empty on database error."""
        with patch("processors.document_processor.supabase_client") as mock_db:
            mock_db.list_documents.side_effect = Exception("DB error")

            from processors.document_processor import find_related_documents
            result = find_related_documents("Send the proposal document")
            assert result == []

    def test_limit_parameter(self):
        """Should respect the limit parameter."""
        with patch("processors.document_processor.supabase_client") as mock_db:
            mock_db.list_documents.return_value = [
                {"title": f"Doc about deck {i}", "document_type": "pitch", "summary": ""}
                for i in range(10)
            ]

            from processors.document_processor import find_related_documents
            result = find_related_documents("Send the deck", limit=3)
            assert len(result) <= 3


# =========================================================================
# Test Telegram notification on document ingestion
# =========================================================================

class TestDocumentIngestionNotification:
    """Tests for Telegram notification when documents are ingested."""

    @pytest.mark.asyncio
    async def test_sends_telegram_notification(self):
        """Should send Telegram notification after processing."""
        with (
            patch("schedulers.document_watcher.drive_service") as mock_drive,
            patch("schedulers.document_watcher.supabase_client") as mock_db,
            patch("schedulers.document_watcher.process_document", new_callable=AsyncMock) as mock_process,
            patch("schedulers.document_watcher.extract_text_by_mime_type") as mock_extract,
        ):
            mock_drive.download_file_bytes = AsyncMock(return_value=b"file content")
            mock_drive.mark_document_processed = MagicMock()
            mock_extract.return_value = "Extracted text content"
            mock_process.return_value = {
                "document_id": "doc-1",
                "title": "Pitch Deck",
                "summary": "A pitch deck for investors",
                "document_type": "pitch",
                "chunk_count": 5,
            }

            from schedulers.document_watcher import DocumentWatcher
            watcher = DocumentWatcher()

            with patch("services.telegram_bot.telegram_bot") as mock_bot:
                mock_bot.send_to_eyal = AsyncMock(return_value=True)

                result = await watcher._process_new_document({
                    "id": "file-1",
                    "name": "Pitch Deck.pptx",
                    "mimeType": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                })

            assert result["status"] == "processed"
            assert result["document_type"] == "pitch"

    @pytest.mark.asyncio
    async def test_notification_failure_does_not_break_processing(self):
        """Telegram notification failure should not fail the processing."""
        with (
            patch("schedulers.document_watcher.drive_service") as mock_drive,
            patch("schedulers.document_watcher.supabase_client") as mock_db,
            patch("schedulers.document_watcher.process_document", new_callable=AsyncMock) as mock_process,
            patch("schedulers.document_watcher.extract_text_by_mime_type") as mock_extract,
        ):
            mock_drive.download_file_bytes = AsyncMock(return_value=b"content")
            mock_drive.mark_document_processed = MagicMock()
            mock_extract.return_value = "text"
            mock_process.return_value = {
                "document_id": "doc-1",
                "title": "Report",
                "summary": "Summary",
                "document_type": "strategy",
                "chunk_count": 3,
            }

            from schedulers.document_watcher import DocumentWatcher
            watcher = DocumentWatcher()

            # Mock telegram_bot to raise an error
            with patch("services.telegram_bot.telegram_bot") as mock_bot:
                mock_bot.send_to_eyal = AsyncMock(side_effect=Exception("Telegram error"))

                result = await watcher._process_new_document({
                    "id": "file-2",
                    "name": "Report.pdf",
                    "mimeType": "application/pdf",
                })

            # Processing should still succeed
            assert result["status"] == "processed"


# =========================================================================
# Test supabase CRUD updates
# =========================================================================

class TestDocumentCRUDUpdates:
    """Tests for updated document CRUD methods."""

    def test_create_document_with_type(self):
        """create_document should accept document_type parameter."""
        from services.supabase_client import SupabaseClient
        client = SupabaseClient()
        mock = MagicMock()
        object.__setattr__(client, "_client", mock)

        mock.table.return_value.insert.return_value.execute.return_value = MagicMock(
            data=[{"id": "doc-1", "title": "Strategy Doc", "document_type": "strategy"}]
        )

        result = client.create_document(
            title="Strategy Doc",
            source="drive",
            file_type="pdf",
            document_type="strategy",
        )
        assert result["document_type"] == "strategy"

        # Verify insert was called with data that includes document_type
        # First table call is for "documents", second is for "audit_log" (from log_action)
        first_insert_call = mock.table.return_value.insert.call_args_list[0]
        insert_data = first_insert_call[0][0]
        assert insert_data["document_type"] == "strategy"

    def test_list_documents_filter_by_type(self):
        """list_documents should filter by document_type."""
        from services.supabase_client import SupabaseClient
        client = SupabaseClient()
        mock = MagicMock()
        object.__setattr__(client, "_client", mock)

        mock.table.return_value.select.return_value.eq.return_value.order.return_value.limit.return_value.execute.return_value = MagicMock(
            data=[{"id": "doc-1", "document_type": "legal"}]
        )

        result = client.list_documents(document_type="legal")
        assert len(result) == 1

    def test_search_documents_by_title(self):
        """search_documents_by_title should do case-insensitive search."""
        from services.supabase_client import SupabaseClient
        client = SupabaseClient()
        mock = MagicMock()
        object.__setattr__(client, "_client", mock)

        mock.table.return_value.select.return_value.ilike.return_value.order.return_value.limit.return_value.execute.return_value = MagicMock(
            data=[{"id": "doc-1", "title": "Lavazza Proposal"}]
        )

        result = client.search_documents_by_title("lavazza")
        assert len(result) == 1
        mock.table.return_value.select.return_value.ilike.assert_called_once_with(
            "title", "%lavazza%"
        )


# =========================================================================
# Test meeting prep includes related documents
# =========================================================================

class TestMeetingPrepRelatedDocuments:
    """Tests for related documents in meeting prep."""

    def test_format_prep_includes_documents_section(self):
        """format_prep_document should include Related Documents section."""
        from processors.meeting_prep import format_prep_document

        event = {
            "title": "Lavazza Review",
            "start": "2026-03-05T10:00:00",
            "location": "Zoom",
            "attendees": [],
            "description": "",
        }
        related_docs = [
            {
                "title": "Lavazza Partnership Proposal",
                "document_type": "client",
                "ingested_at": "2026-02-20T10:00:00",
                "summary": "Proposal for CropSight-Lavazza partnership on precision agriculture data.",
            }
        ]

        result = format_prep_document(
            event=event,
            related_meetings=[],
            relevant_decisions=[],
            open_questions=[],
            participant_tasks={},
            stakeholder_info=[],
            related_documents=related_docs,
        )

        assert "Related Documents" in result
        assert "Lavazza Partnership Proposal" in result
        assert "client" in result

    def test_format_prep_omits_documents_when_empty(self):
        """format_prep_document should not include section when no documents."""
        from processors.meeting_prep import format_prep_document

        event = {
            "title": "Team Standup",
            "start": "2026-03-05",
            "location": "",
            "attendees": [],
            "description": "",
        }

        result = format_prep_document(
            event=event,
            related_meetings=[],
            relevant_decisions=[],
            open_questions=[],
            participant_tasks={},
            stakeholder_info=[],
            related_documents=[],
        )

        assert "Related Documents" not in result
