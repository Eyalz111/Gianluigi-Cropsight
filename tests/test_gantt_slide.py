"""Tests for Sub-Phase 6.4: Gantt slide (PPTX) generation."""

import pytest
from unittest.mock import AsyncMock, MagicMock, patch
from io import BytesIO


@pytest.fixture
def sample_gantt_data():
    return {
        "sections": {
            "Strategic Milestones": {
                "items": [
                    {"name": "★ MVP Release", "owner": "Eyal", "weeks": {"12": "active", "13": "planned"}},
                ]
            },
            "Product & Tech": {
                "items": [
                    {"name": "Yield Model v2", "owner": "Roye", "weeks": {"12": "active", "13": "active", "14": "planned"}},
                    {"name": "API Gateway", "owner": "Roye", "weeks": {"14": "planned", "15": "planned"}},
                ]
            },
            "Sales & BD": {
                "items": [
                    {"name": "● Lavazza PoC", "owner": "Paolo", "weeks": {"12": "active"}},
                ]
            },
            "Fundraising": {
                "items": [
                    {"name": "◆ Seed Round", "owner": "Eyal", "weeks": {"15": "planned", "16": "planned"}},
                ]
            },
            "Legal & Finance": {
                "items": []
            },
        }
    }


# =========================================================================
# Generation Tests
# =========================================================================

class TestGenerateGanttSlide:
    """Test generate_gantt_slide."""

    @pytest.mark.asyncio
    async def test_generates_pptx_bytes(self, sample_gantt_data):
        mock_gm = MagicMock()
        mock_gm.get_gantt_status = AsyncMock(return_value=sample_gantt_data)

        with patch.dict("sys.modules", {"services.gantt_manager": MagicMock(gantt_manager=mock_gm)}):
            from processors.gantt_slide import generate_gantt_slide
            result = await generate_gantt_slide(12, 2026)

        assert isinstance(result, bytes)
        assert len(result) > 0

    @pytest.mark.asyncio
    async def test_pptx_parseable(self, sample_gantt_data):
        """Generated PPTX should be parseable by python-pptx."""
        from pptx import Presentation

        mock_gm = MagicMock()
        mock_gm.get_gantt_status = AsyncMock(return_value=sample_gantt_data)

        with patch.dict("sys.modules", {"services.gantt_manager": MagicMock(gantt_manager=mock_gm)}):
            from processors.gantt_slide import generate_gantt_slide
            pptx_bytes = await generate_gantt_slide(12, 2026)

        prs = Presentation(BytesIO(pptx_bytes))
        assert len(prs.slides) == 1

    @pytest.mark.asyncio
    async def test_empty_gantt_data(self):
        """Should generate slide even with empty Gantt data."""
        mock_gm = MagicMock()
        mock_gm.get_gantt_status = AsyncMock(return_value={})

        with patch.dict("sys.modules", {"services.gantt_manager": MagicMock(gantt_manager=mock_gm)}):
            from processors.gantt_slide import generate_gantt_slide
            result = await generate_gantt_slide(12, 2026)

        assert isinstance(result, bytes)
        assert len(result) > 0

    @pytest.mark.asyncio
    async def test_gantt_unavailable(self):
        """Should generate slide even if Gantt manager fails."""
        with patch.dict("sys.modules", {"services.gantt_manager": None}):
            from processors.gantt_slide import generate_gantt_slide
            result = await generate_gantt_slide(12, 2026)
        assert isinstance(result, bytes)

    @pytest.mark.asyncio
    async def test_custom_week_range(self, sample_gantt_data):
        mock_gm = MagicMock()
        mock_gm.get_gantt_status = AsyncMock(return_value=sample_gantt_data)

        with patch.dict("sys.modules", {"services.gantt_manager": MagicMock(gantt_manager=mock_gm)}):
            from processors.gantt_slide import generate_gantt_slide
            result = await generate_gantt_slide(12, 2026, week_range=(10, 20))
        assert isinstance(result, bytes)


# =========================================================================
# Color Mapping Tests
# =========================================================================

class TestStatusToColor:
    """Test _status_to_color mapping."""

    def test_active(self):
        from processors.gantt_slide import _status_to_color
        from pptx.dml.color import RGBColor
        color = _status_to_color("active")
        assert color == RGBColor(0x4C, 0xAF, 0x50)

    def test_planned(self):
        from processors.gantt_slide import _status_to_color
        color = _status_to_color("planned")
        assert color is not None

    def test_blocked(self):
        from processors.gantt_slide import _status_to_color
        from pptx.dml.color import RGBColor
        color = _status_to_color("blocked")
        assert color == RGBColor(0xF4, 0x43, 0x36)

    def test_completed(self):
        from processors.gantt_slide import _status_to_color
        from pptx.dml.color import RGBColor
        color = _status_to_color("completed")
        assert color == RGBColor(0x9E, 0x9E, 0x9E)

    def test_empty_status(self):
        from processors.gantt_slide import _status_to_color
        color = _status_to_color("")
        assert color is not None

    def test_unknown_status(self):
        from processors.gantt_slide import _status_to_color
        color = _status_to_color("unknown_status")
        assert color is not None


# =========================================================================
# Milestone Marker Tests
# =========================================================================

class TestParseMarkers:
    """Test _parse_markers."""

    def test_tech_marker(self):
        from processors.gantt_slide import _parse_markers
        result = _parse_markers("★ MVP Release")
        assert result["is_milestone"] is True
        assert result["type"] == "tech"

    def test_commercial_marker(self):
        from processors.gantt_slide import _parse_markers
        result = _parse_markers("● First Client")
        assert result["is_milestone"] is True
        assert result["type"] == "commercial"

    def test_funding_marker(self):
        from processors.gantt_slide import _parse_markers
        result = _parse_markers("◆ Seed Round")
        assert result["is_milestone"] is True
        assert result["type"] == "funding"

    def test_no_marker(self):
        from processors.gantt_slide import _parse_markers
        result = _parse_markers("Regular task")
        assert result["is_milestone"] is False


# =========================================================================
# Unicode Handling Tests
# =========================================================================

class TestUnicodeHandling:
    """Test non-Latin text handling in PPTX."""

    @pytest.mark.asyncio
    async def test_hebrew_labels(self):
        """Should handle Hebrew text in Gantt labels."""
        gantt_data = {
            "sections": {
                "Product & Tech": {
                    "items": [
                        {"name": "מודל תשואה", "owner": "Roye", "weeks": {"12": "active"}},
                    ]
                }
            }
        }

        mock_gm = MagicMock()
        mock_gm.get_gantt_status = AsyncMock(return_value=gantt_data)

        with patch.dict("sys.modules", {"services.gantt_manager": MagicMock(gantt_manager=mock_gm)}):
            from processors.gantt_slide import generate_gantt_slide
            result = await generate_gantt_slide(12, 2026)

        assert isinstance(result, bytes)
        assert len(result) > 0

    @pytest.mark.asyncio
    async def test_milestone_markers_unicode(self):
        """★●◆ markers should render correctly."""
        gantt_data = {
            "sections": {
                "Strategic Milestones": {
                    "items": [
                        {"name": "★ Tech milestone", "owner": "Eyal", "weeks": {"12": "active"}},
                        {"name": "● Commercial milestone", "owner": "Paolo", "weeks": {"13": "planned"}},
                        {"name": "◆ Funding milestone", "owner": "Eyal", "weeks": {"14": "planned"}},
                    ]
                }
            }
        }

        mock_gm = MagicMock()
        mock_gm.get_gantt_status = AsyncMock(return_value=gantt_data)

        with patch.dict("sys.modules", {"services.gantt_manager": MagicMock(gantt_manager=mock_gm)}):
            from processors.gantt_slide import generate_gantt_slide
            result = await generate_gantt_slide(12, 2026)

        from pptx import Presentation
        prs = Presentation(BytesIO(result))
        assert len(prs.slides) == 1
